import os
import sys
import tempfile
import gc

_is_windows = sys.platform.startswith("win")
_win32com_available = False

if _is_windows:
    try:
        import win32com.client

        _win32com_available = True
    except ImportError:
        print("Warning: pywin32 is not available. Falling back to python-pptx.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    _python_pptx_available = True
except ImportError:
    _python_pptx_available = False
    if not _win32com_available:
        print("Warning: python-pptx is not available. PowerPoint export is unavailable.")


def add_slide_with_qpixmap(
    ppt_path,
    slide_title,
    slide_text,
    pixmap_images,
    image_positions=None,
    comments_text="",
):
    if not pixmap_images:
        raise ValueError("At least one QPixmap object must be provided.")
    if image_positions is None and len(pixmap_images) != 1:
        raise ValueError("When image_positions is omitted, provide exactly one QPixmap.")
    if image_positions is not None and len(image_positions) != len(pixmap_images):
        raise ValueError("image_positions must match pixmap_images length.")

    temp_files = []

    try:
        for i, pixmap in enumerate(pixmap_images):
            temp_path = os.path.join(
                tempfile.gettempdir(), f"temp_image_{os.getpid()}_{i}.png"
            )
            if not pixmap.save(temp_path, "PNG"):
                raise RuntimeError("Failed to convert QPixmap to PNG before adding to slide.")
            temp_files.append(temp_path)

        if _win32com_available:
            _add_slide_win32com(
                ppt_path,
                slide_title,
                slide_text,
                pixmap_images,
                temp_files,
                image_positions,
                comments_text,
            )
        elif _python_pptx_available:
            _add_slide_python_pptx(
                ppt_path,
                slide_title,
                slide_text,
                pixmap_images,
                temp_files,
                image_positions,
                comments_text,
            )
        else:
            raise RuntimeError(
                "PowerPoint export requires pywin32 on Windows or python-pptx on macOS/Linux."
            )
    finally:
        for temp_path in temp_files:
            try:
                os.remove(temp_path)
            except OSError:
                pass
        gc.collect()


def _add_slide_win32com(
    ppt_path,
    slide_title,
    slide_text,
    pixmap_images,
    temp_files,
    image_positions,
    comments_text,
):
    presentation = None
    ppt_app = None

    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        if not os.path.exists(ppt_path):
            presentation = ppt_app.Presentations.Add()
            presentation.SaveAs(ppt_path)
        else:
            for pres in ppt_app.Presentations:
                if os.path.normcase(pres.FullName) == os.path.normcase(ppt_path):
                    presentation = pres
                    break
            if presentation is None:
                presentation = ppt_app.Presentations.Open(
                    ppt_path, ReadOnly=False, Untitled=False, WithWindow=True
                )

        num_slides = presentation.Slides.Count
        blank_layout = 12
        slide = presentation.Slides.Add(num_slides + 1, blank_layout)

        slide_width = presentation.PageSetup.SlideWidth
        slide_height = presentation.PageSetup.SlideHeight

        # Header area
        title_box = slide.Shapes.AddTextbox(
            Orientation=1, Left=20, Top=10, Width=slide_width - 40, Height=24
        )
        title_box.TextFrame.TextRange.Text = slide_title
        title_box.TextFrame.TextRange.Font.Bold = True
        title_box.TextFrame.TextRange.Font.Size = 16

        meta_box = slide.Shapes.AddTextbox(
            Orientation=1, Left=20, Top=34, Width=slide_width - 40, Height=18
        )
        meta_box.TextFrame.TextRange.Text = slide_text
        meta_box.TextFrame.TextRange.Font.Size = 11

        if image_positions is None:
            # Single screenshot layout: large image left, comments panel right.
            margin = 20
            content_top = 60
            comments_width = max(200, slide_width * 0.28)
            image_box_width = slide_width - comments_width - (3 * margin)
            image_box_height = slide_height - content_top - margin
            image_box_left = margin
            image_box_top = content_top

            pixmap = pixmap_images[0]
            src_w = max(float(pixmap.width()), 1.0)
            src_h = max(float(pixmap.height()), 1.0)
            scale = min(image_box_width / src_w, image_box_height / src_h)

            draw_w = src_w * scale
            draw_h = src_h * scale
            draw_left = image_box_left + (image_box_width - draw_w) / 2
            draw_top = image_box_top + (image_box_height - draw_h) / 2

            slide.Shapes.AddPicture(
                FileName=temp_files[0],
                LinkToFile=False,
                SaveWithDocument=True,
                Left=draw_left,
                Top=draw_top,
                Width=draw_w,
                Height=draw_h,
            )

            comments_left = image_box_left + image_box_width + margin
            comments_height = slide_height - content_top - margin
            comments_value = comments_text.strip() if comments_text and comments_text.strip() else "No comments."
            comments_box = slide.Shapes.AddTextbox(
                Orientation=1,
                Left=comments_left,
                Top=content_top,
                Width=comments_width,
                Height=comments_height,
            )
            comments_box.TextFrame.WordWrap = True
            comments_box.TextFrame.TextRange.Text = f"Comments:\n{comments_value}"
            comments_box.TextFrame.TextRange.Font.Size = 11
        else:
            for temp_path, pixmap, pos in zip(temp_files, pixmap_images, image_positions):
                left, top, width, height = pos
                src_w = max(float(pixmap.width()), 1.0)
                src_h = max(float(pixmap.height()), 1.0)
                scale = min(width / src_w, height / src_h)
                draw_w = src_w * scale
                draw_h = src_h * scale
                draw_left = left + (width - draw_w) / 2
                draw_top = top + (height - draw_h) / 2
                slide.Shapes.AddPicture(
                    FileName=temp_path,
                    LinkToFile=False,
                    SaveWithDocument=True,
                    Left=draw_left,
                    Top=draw_top,
                    Width=draw_w,
                    Height=draw_h,
                )
            if comments_text and comments_text.strip():
                comments_box = slide.Shapes.AddTextbox(
                    Orientation=1,
                    Left=20,
                    Top=slide_height - 70,
                    Width=slide_width - 40,
                    Height=50,
                )
                comments_box.TextFrame.WordWrap = True
                comments_box.TextFrame.TextRange.Text = f"Comments: {comments_text.strip()}"
                comments_box.TextFrame.TextRange.Font.Size = 11

        presentation.Save()
    finally:
        if presentation is not None:
            del presentation
        if ppt_app is not None:
            del ppt_app
        gc.collect()


def _points_to_inches(value):
    return Inches(float(value) / 72.0)


def _fit_rect(source_width, source_height, box_left, box_top, box_width, box_height):
    src_w = max(float(source_width), 1.0)
    src_h = max(float(source_height), 1.0)
    scale = min(float(box_width) / src_w, float(box_height) / src_h)
    draw_w = src_w * scale
    draw_h = src_h * scale
    draw_left = float(box_left) + (float(box_width) - draw_w) / 2
    draw_top = float(box_top) + (float(box_height) - draw_h) / 2
    return draw_left, draw_top, draw_w, draw_h


def _add_slide_python_pptx(
    ppt_path,
    slide_title,
    slide_text,
    pixmap_images,
    temp_files,
    image_positions,
    comments_text,
):
    if os.path.exists(ppt_path):
        presentation = Presentation(ppt_path)
    else:
        presentation = Presentation()

    blank_slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_slide_layout)

    slide_width = presentation.slide_width / 12700.0
    slide_height = presentation.slide_height / 12700.0

    title_box = slide.shapes.add_textbox(
        _points_to_inches(20), _points_to_inches(10), _points_to_inches(slide_width - 40), _points_to_inches(24)
    )
    title_box.text_frame.text = slide_title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(16)

    meta_box = slide.shapes.add_textbox(
        _points_to_inches(20), _points_to_inches(34), _points_to_inches(slide_width - 40), _points_to_inches(18)
    )
    meta_box.text_frame.text = slide_text
    meta_box.text_frame.paragraphs[0].font.size = Pt(11)

    if image_positions is None:
        margin = 20
        content_top = 60
        comments_width = max(200, slide_width * 0.28)
        image_box_width = slide_width - comments_width - (3 * margin)
        image_box_height = slide_height - content_top - margin
        draw_left, draw_top, draw_w, draw_h = _fit_rect(
            pixmap_images[0].width(),
            pixmap_images[0].height(),
            margin,
            content_top,
            image_box_width,
            image_box_height,
        )
        slide.shapes.add_picture(
            temp_files[0],
            _points_to_inches(draw_left),
            _points_to_inches(draw_top),
            width=_points_to_inches(draw_w),
            height=_points_to_inches(draw_h),
        )

        comments_value = (
            comments_text.strip() if comments_text and comments_text.strip() else "No comments."
        )
        comments_box = slide.shapes.add_textbox(
            _points_to_inches(margin + image_box_width + margin),
            _points_to_inches(content_top),
            _points_to_inches(comments_width),
            _points_to_inches(slide_height - content_top - margin),
        )
        comments_box.text_frame.word_wrap = True
        comments_box.text_frame.text = f"Comments:\n{comments_value}"
        comments_box.text_frame.paragraphs[0].font.size = Pt(11)
    else:
        for temp_path, pixmap, pos in zip(temp_files, pixmap_images, image_positions):
            left, top, width, height = pos
            draw_left, draw_top, draw_w, draw_h = _fit_rect(
                pixmap.width(), pixmap.height(), left, top, width, height
            )
            slide.shapes.add_picture(
                temp_path,
                _points_to_inches(draw_left),
                _points_to_inches(draw_top),
                width=_points_to_inches(draw_w),
                height=_points_to_inches(draw_h),
            )

        if comments_text and comments_text.strip():
            comments_box = slide.shapes.add_textbox(
                _points_to_inches(20),
                _points_to_inches(slide_height - 70),
                _points_to_inches(slide_width - 40),
                _points_to_inches(50),
            )
            comments_box.text_frame.word_wrap = True
            comments_box.text_frame.text = f"Comments: {comments_text.strip()}"
            comments_box.text_frame.paragraphs[0].font.size = Pt(11)

    presentation.save(ppt_path)


if __name__ == "__main__":
    from PyQt6.QtWidgets import QApplication, QWidget

    app = QApplication(sys.argv)
    widget = QWidget()
    widget.resize(300, 200)
    widget.show()
    app.processEvents()

    # Grab images from the widget
    image1 = widget.grab()
    image2 = widget.grab()
    image3 = widget.grab()

    ppt_file = r"C:\Users\moham\OneDrive\Documents\test_presentation.pptx"
    title = "Slide with QPixmap Images"
    text = "This slide contains three images grabbed from PyQt widgets."

    add_slide_with_qpixmap(
        ppt_file,
        title,
        text,
        [image1, image2, image3],
        image_positions=[(20, 90, 320, 220), (360, 90, 320, 220), (20, 330, 320, 180)],
    )
    
    # Do not call app.exec_() to avoid blocking; close the application instead.
    widget.close()
    app.exit()
